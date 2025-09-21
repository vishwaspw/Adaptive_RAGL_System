import streamlit as st
import os
import json
from datetime import datetime
import time
from typing import Dict, List, Any, Optional
import sys

# Add the parent directory to path to resolve imports
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from src.components.ragl_system import RAGLSystem
from src.utils.config import GEMINI_API_KEY

# Page configuration
st.set_page_config(
    page_title="Adaptive RAGL Chatbot",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="expanded",
)

# Custom CSS for better UI
st.markdown("""
<style>
.chat-message {
    padding: 1.5rem; border-radius: 0.5rem; margin-bottom: 1rem; display: flex;
    background-color: #f0f2f6;
    color: #0e1117;
}
.chat-message.user {
    background-color: #f0f2f6;
    color: #0e1117;
}
.chat-message.assistant {
    background-color: #d1e8ff;
    color: #0e1117;
}
.chat-message .avatar {
    width: 40px; margin-right: 1rem;
}
.chat-message .message {
    flex-grow: 1;
    color: #0e1117;
}
.learning-banner {
    padding: 0.5rem; 
    border-radius: 0.5rem; 
    background-color: #ffeb99;
    color: #0e1117;
    margin-bottom: 1rem;
    text-align: center;
    font-weight: 600;
}
.document-box {
    border: 1px solid #ddd;
    padding: 0.8rem;
    border-radius: 0.3rem;
    margin-bottom: 0.5rem;
    background-color: #f9f9f9;
    color: #0e1117;
}
.footer {
    text-align: center;
    margin-top: 2rem;
    font-size: 0.8rem;
    color: #b0b0b0;
}
/* Improve text visibility in the dark theme */
.st-emotion-cache-ue6h4q {
    color: #ffffff;
}
.st-emotion-cache-ue6h4q a {
    color: #4da6ff;
}
/* Make sure input text is visible */
.stTextArea textarea {
    color: #0e1117;
    background-color: #ffffff;
}
</style>
""", unsafe_allow_html=True)

# Initialize session state
if "ragl_system" not in st.session_state:
    st.session_state.ragl_system = RAGLSystem()

if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

if "learning_events" not in st.session_state:
    st.session_state.learning_events = []

if "document_feedback" not in st.session_state:
    st.session_state.document_feedback = {}

# App title
st.title("🧠 Adaptive RAGL Chatbot")

# Sidebar
with st.sidebar:
    st.header("About")
    st.markdown("""
    This chatbot is powered by an Adaptive Retrieval-Augmented Generation Learning (RAGL) system
    that continuously curates and updates its knowledge base based on interactions.
    """)
    
    st.header("Settings")
    
    # Document upload
    st.subheader("Upload Documents")
    uploaded_files = st.file_uploader(
        "Add documents to the knowledge base", 
        accept_multiple_files=True, 
        type=["txt", "md", "pdf", "csv", "doc", "docx", "xls", "xlsx", "ppt", "pptx", "json", "xml", "html", "htm"]
    )
    
    if uploaded_files:
        if st.button("Process Uploaded Documents"):
            with st.spinner("Processing documents..."):
                documents = []
                metadatas = []
                
                for file in uploaded_files:
                    try:
                        # Read file content based on file type
                        file_bytes = file.read()
                        content = ""
                        
                        # Handle different file types
                        if file.type == "application/pdf":
                            try:
                                import PyPDF2
                                from io import BytesIO
                                # Convert bytes to file-like object
                                pdf_file = BytesIO(file_bytes)
                                pdf_reader = PyPDF2.PdfReader(pdf_file)
                                content = ""
                                for page in pdf_reader.pages:
                                    content += page.extract_text() + "\n"
                            except Exception as e:
                                st.warning(f"Could not extract text from PDF {file.name}: {str(e)}")
                                content = f"[PDF content from file: {file.name}]"
                        
                        elif file.type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
                                         "application/msword"]:
                            try:
                                import docx
                                doc = docx.Document(file_bytes)
                                content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                            except Exception as e:
                                st.warning(f"Could not extract text from Word document {file.name}: {str(e)}")
                                content = f"[Word document content from file: {file.name}]"
                        
                        elif file.type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                         "application/vnd.ms-excel"]:
                            try:
                                import pandas as pd
                                df = pd.read_excel(file_bytes)
                                content = df.to_string()
                            except Exception as e:
                                st.warning(f"Could not extract text from Excel file {file.name}: {str(e)}")
                                content = f"[Excel content from file: {file.name}]"
                        
                        elif file.type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                         "application/vnd.ms-powerpoint"]:
                            try:
                                from pptx import Presentation
                                prs = Presentation(file_bytes)
                                content = ""
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if hasattr(shape, "text"):
                                            content += shape.text + "\n"
                            except Exception as e:
                                st.warning(f"Could not extract text from PowerPoint file {file.name}: {str(e)}")
                                content = f"[PowerPoint content from file: {file.name}]"
                        
                        elif file.type == "text/csv":
                            try:
                                import pandas as pd
                                df = pd.read_csv(file_bytes)
                                content = df.to_string()
                            except Exception as e:
                                st.warning(f"Could not extract text from CSV file {file.name}: {str(e)}")
                                content = f"[CSV content from file: {file.name}]"
                        
                        elif file.type in ["application/json", "text/xml", "text/html"]:
                            try:
                                content = file_bytes.decode("utf-8")
                            except UnicodeDecodeError:
                                content = f"[Structured content from file: {file.name}]"
                        
                        else:
                            # Try to decode as text for other file types
                            try:
                                content = file_bytes.decode("utf-8")
                            except UnicodeDecodeError:
                                content = f"[Binary content from file: {file.name}]"
                        
                        # Create metadata
                        metadata = {
                            "filename": file.name,
                            "file_type": file.type,
                            "source": "user_upload",
                            "upload_timestamp": datetime.now().isoformat(),
                            "content_type": "text" if isinstance(content, str) else "binary"
                        }
                        
                        documents.append(content)
                        metadatas.append(metadata)
                        
                        # Show success message for each file
                        st.success(f"Successfully processed {file.name}")
                        
                    except Exception as e:
                        st.error(f"Error processing file {file.name}: {str(e)}")
                
                if documents:
                    # Add to RAGL system
                    doc_ids = st.session_state.ragl_system.add_documents(documents, metadatas)
                    st.success(f"Added {len(doc_ids)} documents to the knowledge base!")
                    
                    # Show summary of added documents
                    with st.expander("View Added Documents Summary"):
                        for i, (doc, meta) in enumerate(zip(documents, metadatas)):
                            st.markdown(f"**Document {i+1}:** {meta['filename']}")
                            st.markdown(f"Type: {meta['file_type']}")
                            st.markdown(f"Content Preview: {doc[:200]}...")
                            st.markdown("---")
    
    # View learning events
    st.subheader("Learning Events")
    if st.button("View Learning Events"):
        learning_events = st.session_state.learning_events
        if learning_events:
            for i, event in enumerate(learning_events):
                st.markdown(f"**Event {i+1}**")
                st.markdown(f"**Knowledge:** {event['knowledge']}")
                st.markdown(f"**Confidence:** {event['confidence']:.2f}")
                st.markdown(f"**Time:** {event['timestamp']}")
                st.markdown("---")
        else:
            st.info("No learning events recorded yet.")
    
    # Save session
    st.subheader("Save Session")
    if st.button("Save Current Session"):
        file_path = st.session_state.ragl_system.save_session()
        st.success(f"Session saved to {file_path}")
    
    # API key status indicator
    st.subheader("API Key Status")
    if GEMINI_API_KEY and GEMINI_API_KEY != "your_gemini_api_key_here":
        st.success("✅ Gemini API key configured")
    else:
        st.error("❌ Gemini API key not configured. Please add it to the .env file.")

# Main chat interface
chat_container = st.container()

# Document display container
doc_container = st.container()

# Function to display chat messages
def display_chat():
    with chat_container:
        # Clear previous messages before displaying updated chat
        chat_container.empty()
        
        for message in st.session_state.chat_history:
            role = message["role"]
            content = message["content"]
            
            if "learning_occurred" in message and message["learning_occurred"]:
                st.markdown('<div class="learning-banner">🧠 New knowledge acquired and added to the system!</div>', unsafe_allow_html=True)
            
            if role == "user":
                st.markdown(f'<div class="chat-message user"><div class="avatar">👤</div><div class="message">{content}</div></div>', unsafe_allow_html=True)
            else:
                st.markdown(f'<div class="chat-message assistant"><div class="avatar">🤖</div><div class="message">{content}</div></div>', unsafe_allow_html=True)

# Function to display retrieved documents
def display_documents(documents):
    with doc_container:
        # Clear previous documents before displaying
        doc_container.empty()
        
        if documents:
            with st.expander("View Retrieved Documents", expanded=False):
                for i, doc in enumerate(documents):
                    st.markdown(f'<div class="document-box">', unsafe_allow_html=True)
                    st.markdown(f"**Document {i+1}**")
                    st.markdown(doc["content"])
                    
                    # Metadata display
                    if "metadata" in doc and doc["metadata"]:
                        meta_cols = st.columns(2)
                        with meta_cols[0]:
                            if "source" in doc["metadata"]:
                                st.caption(f"Source: {doc['metadata']['source']}")
                            if "timestamp" in doc["metadata"]:
                                st.caption(f"Added: {doc['metadata']['timestamp'][:10]}")
                        
                        with meta_cols[1]:
                            # Relevance feedback
                            doc_key = str(i)  # Use document index as key
                            if doc_key not in st.session_state.document_feedback:
                                st.session_state.document_feedback[doc_key] = 0.5
                            
                            st.caption("Document relevance:")
                            feedback = st.slider("", 0.0, 1.0, st.session_state.document_feedback[doc_key], key=f"feedback_{i}")
                            
                            if feedback != st.session_state.document_feedback[doc_key]:
                                st.session_state.document_feedback[doc_key] = feedback
                                # Store feedback to process later
                    
                    st.markdown('</div>', unsafe_allow_html=True)
        else:
            st.info("No relevant documents were retrieved for this query.")

# Process user input
with st.form(key="user_input_form", clear_on_submit=True):
    user_input = st.text_area("Type your message:", placeholder="Ask me anything...", height=100)
    submit_button = st.form_submit_button("Send")

# Process the message when user hits send
if submit_button and user_input:
    # Add user message to chat history
    st.session_state.chat_history.append({"role": "user", "content": user_input})
    
    # Display updated chat immediately (just the user's message)
    display_chat()
    
    # Get response from RAGL system
    with st.spinner("Thinking..."):
        try:
            # Create a placeholder for the assistant's message that will be updated
            message_placeholder = st.empty()
            
            result = st.session_state.ragl_system.process_query(user_input)
            response = result["response"]
            retrieved_docs = result["retrieved_docs"]
            learning_occurred = result["learning_occurred"]
            
            # Process document feedback from previous interaction
            for doc_key, score in st.session_state.document_feedback.items():
                try:
                    doc_id = int(doc_key)
                    # Ideally we'd map this to actual document IDs
                    st.session_state.ragl_system.provide_feedback(doc_id, score)
                except:
                    pass
            
            # Clear document feedback for new interaction
            st.session_state.document_feedback = {}
            
            # Add assistant message to chat history
            st.session_state.chat_history.append({
                "role": "assistant", 
                "content": response,
                "learning_occurred": learning_occurred
            })
            
            # If learning occurred, add to learning events
            if learning_occurred and st.session_state.ragl_system.learning_events:
                latest_event = st.session_state.ragl_system.learning_events[-1]
                st.session_state.learning_events.append(latest_event)
            
            # Display updated chat with a slight delay to ensure proper rendering
            st.rerun()
            
        except Exception as e:
            st.error(f"Error: {str(e)}")
            if "API key" in str(e):
                st.error("Please configure your Gemini API key in the .env file.")
else:
    # Just display the existing chat
    display_chat()

# Footer
st.markdown("""
<div class="footer">
    Adaptive RAGL Chatbot | Powered by Google Gemini & LangChain
</div>
""", unsafe_allow_html=True)

# Initial instructions
if not st.session_state.chat_history:
    st.info("""
    👋 Welcome to the Adaptive RAGL Chatbot!
    
    This system will:
    - Answer your questions using its knowledge base
    - Learn from your interactions
    - Continuously improve its knowledge
    
    Try asking a question or sharing some information!
    """)

# Update script when API key is missing
if not GEMINI_API_KEY or GEMINI_API_KEY == "your_gemini_api_key_here":
    st.warning("""
    ⚠️ **API Key Not Configured**
    
    Please set your Gemini API key in the `.env` file:
    1. Open the `.env` file in the project root
    2. Replace `your_gemini_api_key_here` with your actual API key
    3. Save the file and restart the application
    """)

# Add an __init__.py file to each directory to make the imports work
if not os.path.exists("src/__init__.py"):
    with open("src/__init__.py", "w") as f:
        pass

if not os.path.exists("src/utils/__init__.py"):
    with open("src/utils/__init__.py", "w") as f:
        pass

if not os.path.exists("src/components/__init__.py"):
    with open("src/components/__init__.py", "w") as f:
        pass 